"""Utilities for generating the Cyber intrusion story presentation without committing binaries.

This script recreates the five-slide stakeholder deck using the ``Cyber_pres.pptx`` template
shipped with the repository. The presentation mirrors the previously supplied
``reports/Cyber_intrusion_story.pptx`` file, but it can be regenerated on demand so
contributors do not need to add large binary assets to version control.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = ROOT / "Cyber_pres.pptx"
DEFAULT_OUTPUT = ROOT / "reports" / "Cyber_intrusion_story.pptx"


def _clear_template_slides(prs: Presentation) -> None:
    """Remove any placeholder slides bundled with the template."""
    while len(prs.slides):
        r_id = prs.slides._sldIdLst[0].rId  # type: ignore[attr-defined]
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[0]  # type: ignore[attr-defined]


def _add_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Project Cyber: Intrusion Detection Dashboard"
    subtitle = slide.placeholders[1]
    subtitle.text = (
        "Entertaining threat intelligence for faster decisions\n"
        "Presented by the Cyber Analytics Guild"
    )


def _add_mission_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Mission Brief"
    body = slide.shapes.placeholders[1].text_frame

    body.text = "Project Scope"
    for bullet in (
        "Deliver a lean intrusion detection dashboard for lightning-fast situational awareness.",
        "Blend machine learning with narrative storytelling to keep security teams engaged.",
        "Expose critical threat patterns, user behaviour shifts, and model transparency in one hub.",
    ):
        p = body.add_paragraph()
        p.text = bullet
        p.level = 1

    audience = body.add_paragraph()
    audience.text = "Primary Audience: SOC leads, CISOs, and data-savvy executives."

    body.add_paragraph().text = "Narrative Hook"
    for hook in (
        "Open with a “Breach of the Week” cold open to humanise the risk.",
        "Colour-coded threat levels and plain-language tooltips demystify model outputs.",
        "Celebrate analyst victories—blocked attacks, shortened dwell times, restored trust.",
        "Every visual answers the stakeholder favourite: “So what?”",
    ):
        p = body.add_paragraph()
        p.text = hook
        p.level = 1

    hypothesis = body.add_paragraph()
    hypothesis.text = "Core Hypothesis: Traffic spikes on risky protocols foreshadow malicious activity."

    body.add_paragraph().text = "Toolkit"
    for item in (
        "Python (Pandas, Scikit-learn) for feature wrangling and precision-first models.",
        "Streamlit app mirrors Tableau prototype for live-fire demos.",
        "GitHub projects orchestrate version control, issues, and release cadence.",
        "Automated profiling + SHAP storytelling keeps the narrative evidence-based.",
        "Alert scripts plug into Slack/email for rapid incident choreography.",
    ):
        p = body.add_paragraph()
        p.text = item
        p.level = 1


def _add_data_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Data & Modelling Sprints"
    body = slide.shapes.placeholders[1].text_frame
    body.text = (
        "Dataset: 9,537 labelled network sessions packed with protocol, port, packet, and behaviour signals."
    )

    for bullet in (
        "Preprocessing: Stratified sampling, null remediation, scaling, and categorical encoding for fairness.",
        "Target Intelligence: attack_detected label guides supervised threat classification.",
        "Validation: Stratified K-fold achieving ≥0.90 precision with ≤5% false positives.",
        "Explainability: SHAP radar spotlights packet burst size, unusual ports, and failed logins.",
        "Deployment: Streamlit dashboard delivers filterable tables, KPI cards, and anomaly callouts.",
    ):
        p = body.add_paragraph()
        p.text = bullet


def _add_insights_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Key Insights & Visual Stories"
    body = slide.shapes.placeholders[1].text_frame
    body.text = (
        "High packet surges on non-standard ports triple intrusion odds—our treemap turns red as the attacks roll in."
    )
    p = body.add_paragraph()
    p.text = (
        "After-hours sessions using weak encryption spark 2.4× more alerts; the timeline heatmap dramatises the witching hour."
    )


def _add_recommendations_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Recommendations & Next Missions"
    body = slide.shapes.placeholders[1].text_frame
    body.text = (
        "Automate alert escalations with risk-tiered Slack pings and ready-to-go analyst playbooks."
    )
    p = body.add_paragraph()
    p.text = (
        "Feed confirmed incidents back into the model for adaptive learning and richer SOC feedback loops."
    )


def build_presentation(output_path: Path = DEFAULT_OUTPUT) -> Path:
    """Create the PowerPoint file at ``output_path``."""
    prs = Presentation(str(TEMPLATE))
    _clear_template_slides(prs)
    _add_title_slide(prs)
    _add_mission_slide(prs)
    _add_data_slide(prs)
    _add_insights_slide(prs)
    _add_recommendations_slide(prs)

    output_path = output_path.resolve()
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path


def main() -> None:
    output = build_presentation()
    print(f"Presentation written to {output}")


if __name__ == "__main__":
    main()
